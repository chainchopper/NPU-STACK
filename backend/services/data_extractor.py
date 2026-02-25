"""Data Extractor — Universal file-to-text/data extraction for 15+ file types.

Extracts structured content from documents, images, audio, and data files
to feed into the dataset builder for training data creation.
"""

import os
import json
import csv
import io
import mimetypes
from typing import Optional


# ── Supported Extensions ────────────────────────────────
SUPPORTED_EXTENSIONS = {
    # Documents
    ".pdf", ".doc", ".docx", ".ppt", ".pptx", ".odt",
    # Text/markup
    ".txt", ".md", ".rst", ".log", ".rtf",
    # Data formats
    ".json", ".jsonl", ".ndjson", ".xml", ".csv", ".tsv", ".parquet", ".yaml", ".yml",
    # Images
    ".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".webp",
    # Audio
    ".wav", ".mp3", ".flac", ".ogg", ".m4a",
    # Web
    ".html", ".htm",
}


def get_supported_types() -> list:
    """Return list of supported file types grouped by category."""
    return [
        {"category": "Documents", "extensions": [".pdf", ".docx", ".doc", ".pptx", ".ppt", ".odt"]},
        {"category": "Text & Markup", "extensions": [".txt", ".md", ".rst", ".log", ".rtf"]},
        {"category": "Data Files", "extensions": [".json", ".jsonl", ".ndjson", ".csv", ".tsv", ".xml", ".parquet", ".yaml", ".yml"]},
        {"category": "Images", "extensions": [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".webp"]},
        {"category": "Audio", "extensions": [".wav", ".mp3", ".flac", ".ogg", ".m4a"]},
        {"category": "Web", "extensions": [".html", ".htm"]},
    ]


def extract_file(file_path: str, ocr: bool = False) -> dict:
    """Extract text/data from a file. Returns structured extraction result.

    Args:
        file_path: Absolute path to the file.
        ocr: If True, attempt OCR on images.

    Returns:
        dict with keys: text, metadata, file_type, chunks, success, error
    """
    if not os.path.exists(file_path):
        return _error(file_path, "File not found")

    ext = os.path.splitext(file_path)[1].lower()
    if ext not in SUPPORTED_EXTENSIONS:
        return _error(file_path, f"Unsupported file type: {ext}")

    file_size = os.path.getsize(file_path)
    base_meta = {
        "file_name": os.path.basename(file_path),
        "file_path": file_path,
        "file_size": file_size,
        "file_size_human": _format_size(file_size),
        "extension": ext,
    }

    try:
        if ext == ".pdf":
            return _extract_pdf(file_path, base_meta, ocr)
        elif ext in (".doc", ".docx"):
            return _extract_docx(file_path, base_meta)
        elif ext in (".ppt", ".pptx"):
            return _extract_pptx(file_path, base_meta)
        elif ext in (".txt", ".md", ".rst", ".log", ".rtf"):
            return _extract_text(file_path, base_meta)
        elif ext in (".json",):
            return _extract_json(file_path, base_meta)
        elif ext in (".jsonl", ".ndjson"):
            return _extract_jsonl(file_path, base_meta)
        elif ext in (".csv", ".tsv"):
            return _extract_csv(file_path, base_meta, delimiter="\t" if ext == ".tsv" else ",")
        elif ext == ".parquet":
            return _extract_parquet(file_path, base_meta)
        elif ext == ".xml":
            return _extract_xml(file_path, base_meta)
        elif ext in (".yaml", ".yml"):
            return _extract_yaml(file_path, base_meta)
        elif ext in (".html", ".htm"):
            return _extract_html(file_path, base_meta)
        elif ext in (".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".webp"):
            return _extract_image(file_path, base_meta, ocr)
        elif ext in (".wav", ".mp3", ".flac", ".ogg", ".m4a"):
            return _extract_audio(file_path, base_meta)
        else:
            return _error(file_path, f"No extractor for {ext}")
    except Exception as e:
        return _error(file_path, str(e), base_meta)


def extract_folder(folder_path: str, recursive: bool = True, ocr: bool = False) -> dict:
    """Extract data from all supported files in a folder.

    Returns:
        dict with keys: files (list of extraction results), total, success_count, error_count
    """
    results = []
    for root, dirs, files in os.walk(folder_path) if recursive else [(folder_path, [], os.listdir(folder_path))]:
        for fname in files:
            ext = os.path.splitext(fname)[1].lower()
            if ext in SUPPORTED_EXTENSIONS:
                fpath = os.path.join(root, fname)
                results.append(extract_file(fpath, ocr=ocr))

    success_count = sum(1 for r in results if r.get("success"))
    return {
        "files": results,
        "total": len(results),
        "success_count": success_count,
        "error_count": len(results) - success_count,
    }


# ── Extractors ──────────────────────────────────────────

def _extract_pdf(path: str, meta: dict, ocr: bool) -> dict:
    try:
        import fitz  # PyMuPDF
    except ImportError:
        try:
            import pdfplumber
            return _extract_pdf_plumber(path, meta)
        except ImportError:
            return _error(path, "Install PyMuPDF or pdfplumber: pip install pymupdf pdfplumber", meta)

    doc = fitz.open(path)
    pages = []
    full_text = []
    for i, page in enumerate(doc):
        text = page.get_text()
        pages.append({"page": i + 1, "text": text, "char_count": len(text)})
        full_text.append(text)
    doc.close()

    text_joined = "\n\n".join(full_text)
    meta["page_count"] = len(pages)
    meta["char_count"] = len(text_joined)
    meta["word_count"] = len(text_joined.split())

    return {
        "success": True,
        "file_type": "pdf",
        "text": text_joined,
        "chunks": pages,
        "metadata": meta,
    }


def _extract_pdf_plumber(path: str, meta: dict) -> dict:
    import pdfplumber
    pages = []
    full_text = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            pages.append({"page": i + 1, "text": text, "char_count": len(text)})
            full_text.append(text)

    text_joined = "\n\n".join(full_text)
    meta["page_count"] = len(pages)
    meta["char_count"] = len(text_joined)

    return {
        "success": True,
        "file_type": "pdf",
        "text": text_joined,
        "chunks": pages,
        "metadata": meta,
    }


def _extract_docx(path: str, meta: dict) -> dict:
    try:
        from docx import Document
    except ImportError:
        return _error(path, "Install python-docx: pip install python-docx", meta)

    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    # Extract tables
    tables = []
    for table in doc.tables:
        rows = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        tables.append(rows)

    text = "\n\n".join(paragraphs)
    meta["paragraph_count"] = len(paragraphs)
    meta["table_count"] = len(tables)
    meta["char_count"] = len(text)

    return {
        "success": True,
        "file_type": "docx",
        "text": text,
        "chunks": [{"type": "paragraph", "text": p} for p in paragraphs],
        "tables": tables,
        "metadata": meta,
    }


def _extract_pptx(path: str, meta: dict) -> dict:
    try:
        from pptx import Presentation
    except ImportError:
        return _error(path, "Install python-pptx: pip install python-pptx", meta)

    prs = Presentation(path)
    slides = []
    full_text = []
    for i, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
        slide_text = "\n".join(slide_texts)
        slides.append({"slide": i + 1, "text": slide_text, "notes": notes})
        full_text.append(slide_text)
        if notes:
            full_text.append(f"[Notes] {notes}")

    text = "\n\n".join(full_text)
    meta["slide_count"] = len(slides)
    meta["char_count"] = len(text)

    return {
        "success": True,
        "file_type": "pptx",
        "text": text,
        "chunks": slides,
        "metadata": meta,
    }


def _extract_text(path: str, meta: dict) -> dict:
    # Try multiple encodings
    for enc in ["utf-8", "utf-8-sig", "latin-1", "cp1252"]:
        try:
            with open(path, "r", encoding=enc) as f:
                text = f.read()
            meta["encoding"] = enc
            meta["char_count"] = len(text)
            meta["line_count"] = text.count("\n") + 1
            meta["word_count"] = len(text.split())

            return {
                "success": True,
                "file_type": "text",
                "text": text,
                "chunks": [{"text": text}],
                "metadata": meta,
            }
        except (UnicodeDecodeError, UnicodeError):
            continue

    return _error(path, "Failed to decode text with any supported encoding", meta)


def _extract_json(path: str, meta: dict) -> dict:
    with open(path, "r", encoding="utf-8") as f:
        data = json.load(f)

    if isinstance(data, list):
        meta["record_count"] = len(data)
        meta["schema"] = _infer_schema(data[0]) if data else {}
        text = json.dumps(data, indent=2, ensure_ascii=False)
        chunks = [{"index": i, "data": item} for i, item in enumerate(data[:100])]
    elif isinstance(data, dict):
        meta["key_count"] = len(data)
        text = json.dumps(data, indent=2, ensure_ascii=False)
        chunks = [{"key": k, "value": v} for k, v in list(data.items())[:100]]
    else:
        text = json.dumps(data, ensure_ascii=False)
        chunks = [{"data": data}]

    meta["char_count"] = len(text)

    return {
        "success": True,
        "file_type": "json",
        "text": text,
        "data": data,
        "chunks": chunks,
        "metadata": meta,
    }


def _extract_jsonl(path: str, meta: dict) -> dict:
    records = []
    with open(path, "r", encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if line:
                try:
                    records.append(json.loads(line))
                except json.JSONDecodeError:
                    continue

    meta["record_count"] = len(records)
    meta["schema"] = _infer_schema(records[0]) if records else {}
    text = "\n".join(json.dumps(r, ensure_ascii=False) for r in records[:50])

    return {
        "success": True,
        "file_type": "jsonl",
        "text": text,
        "data": records[:100],
        "chunks": [{"index": i, "data": r} for i, r in enumerate(records[:100])],
        "metadata": meta,
    }


def _extract_csv(path: str, meta: dict, delimiter: str = ",") -> dict:
    rows = []
    with open(path, "r", encoding="utf-8", newline="") as f:
        reader = csv.DictReader(f, delimiter=delimiter)
        headers = reader.fieldnames or []
        for row in reader:
            rows.append(dict(row))

    meta["row_count"] = len(rows)
    meta["columns"] = headers
    meta["column_count"] = len(headers)

    # Build text representation
    text_lines = [delimiter.join(headers)]
    for row in rows[:50]:
        text_lines.append(delimiter.join(str(v) for v in row.values()))
    text = "\n".join(text_lines)

    return {
        "success": True,
        "file_type": "csv",
        "text": text,
        "data": rows[:100],
        "chunks": [{"index": i, "data": r} for i, r in enumerate(rows[:100])],
        "metadata": meta,
    }


def _extract_parquet(path: str, meta: dict) -> dict:
    try:
        import pandas as pd
    except ImportError:
        return _error(path, "Install pandas + pyarrow: pip install pandas pyarrow", meta)

    df = pd.read_parquet(path)
    meta["row_count"] = len(df)
    meta["columns"] = list(df.columns)
    meta["column_count"] = len(df.columns)
    meta["dtypes"] = {col: str(dtype) for col, dtype in df.dtypes.items()}

    records = df.head(100).to_dict(orient="records")
    text = df.head(50).to_string()

    return {
        "success": True,
        "file_type": "parquet",
        "text": text,
        "data": records,
        "chunks": [{"index": i, "data": r} for i, r in enumerate(records)],
        "metadata": meta,
    }


def _extract_xml(path: str, meta: dict) -> dict:
    import xml.etree.ElementTree as ET

    tree = ET.parse(path)
    root = tree.getroot()

    def _elem_to_text(elem, depth=0):
        parts = []
        tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
        text = (elem.text or "").strip()
        if text:
            parts.append(f"{'  ' * depth}{tag}: {text}")
        for child in elem:
            parts.extend(_elem_to_text(child, depth + 1))
        tail = (elem.tail or "").strip()
        if tail:
            parts.append(f"{'  ' * depth}{tail}")
        return parts

    lines = _elem_to_text(root)
    text = "\n".join(lines)
    meta["root_tag"] = root.tag
    meta["char_count"] = len(text)

    return {
        "success": True,
        "file_type": "xml",
        "text": text,
        "chunks": [{"text": text}],
        "metadata": meta,
    }


def _extract_yaml(path: str, meta: dict) -> dict:
    try:
        import yaml
    except ImportError:
        return _error(path, "Install pyyaml: pip install pyyaml", meta)

    with open(path, "r", encoding="utf-8") as f:
        data = yaml.safe_load(f)

    text = json.dumps(data, indent=2, ensure_ascii=False) if data else ""
    meta["char_count"] = len(text)

    return {
        "success": True,
        "file_type": "yaml",
        "text": text,
        "data": data,
        "chunks": [{"data": data}],
        "metadata": meta,
    }


def _extract_html(path: str, meta: dict) -> dict:
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        # Fallback: strip tags manually
        with open(path, "r", encoding="utf-8") as f:
            import re
            raw = f.read()
            text = re.sub(r"<[^>]+>", " ", raw)
            text = re.sub(r"\s+", " ", text).strip()
        meta["char_count"] = len(text)
        return {"success": True, "file_type": "html", "text": text, "chunks": [{"text": text}], "metadata": meta}

    with open(path, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f, "html.parser")

    title = soup.title.string if soup.title else ""
    text = soup.get_text(separator="\n", strip=True)
    meta["title"] = title
    meta["char_count"] = len(text)

    return {
        "success": True,
        "file_type": "html",
        "text": text,
        "chunks": [{"text": text}],
        "metadata": meta,
    }


def _extract_image(path: str, meta: dict, ocr: bool) -> dict:
    try:
        from PIL import Image
    except ImportError:
        return _error(path, "Install Pillow: pip install Pillow", meta)

    img = Image.open(path)
    meta["width"] = img.width
    meta["height"] = img.height
    meta["mode"] = img.mode
    meta["format"] = img.format

    text = ""
    caption = ""

    # OCR if requested
    if ocr:
        try:
            import pytesseract
            text = pytesseract.image_to_string(img)
            meta["ocr_applied"] = True
        except ImportError:
            meta["ocr_applied"] = False
            meta["ocr_note"] = "Install pytesseract for OCR: pip install pytesseract"

    img.close()

    return {
        "success": True,
        "file_type": "image",
        "text": text,
        "caption": caption,
        "chunks": [{"text": text}] if text else [],
        "metadata": meta,
    }


def _extract_audio(path: str, meta: dict) -> dict:
    # Get audio metadata
    try:
        import librosa
        y, sr = librosa.load(path, sr=None, duration=10)
        duration = librosa.get_duration(y=y, sr=sr)
        meta["sample_rate"] = sr
        meta["duration_seconds"] = round(duration, 2)
    except ImportError:
        try:
            from pydub import AudioSegment
            audio = AudioSegment.from_file(path)
            meta["duration_seconds"] = round(len(audio) / 1000, 2)
            meta["sample_rate"] = audio.frame_rate
            meta["channels"] = audio.channels
        except ImportError:
            meta["duration_seconds"] = None
            meta["note"] = "Install librosa or pydub for audio metadata"

    text = ""
    # Try Whisper transcription
    try:
        import whisper
        model = whisper.load_model("base")
        result = model.transcribe(path)
        text = result.get("text", "")
        meta["transcription_model"] = "whisper-base"
    except ImportError:
        meta["transcription_note"] = "Install openai-whisper for audio transcription"

    return {
        "success": True,
        "file_type": "audio",
        "text": text,
        "chunks": [{"text": text}] if text else [],
        "metadata": meta,
    }


# ── Utilities ───────────────────────────────────────────

def _error(path: str, message: str, meta: dict = None) -> dict:
    return {
        "success": False,
        "file_type": "unknown",
        "text": "",
        "chunks": [],
        "error": message,
        "metadata": meta or {"file_path": path, "file_name": os.path.basename(path)},
    }


def _infer_schema(record: dict) -> dict:
    """Infer field types from a sample record."""
    if not isinstance(record, dict):
        return {}
    return {k: type(v).__name__ for k, v in record.items()}


def _format_size(size_bytes: int) -> str:
    if size_bytes < 1024:
        return f"{size_bytes} B"
    elif size_bytes < 1024 ** 2:
        return f"{size_bytes / 1024:.1f} KB"
    elif size_bytes < 1024 ** 3:
        return f"{size_bytes / (1024 ** 2):.1f} MB"
    else:
        return f"{size_bytes / (1024 ** 3):.2f} GB"
